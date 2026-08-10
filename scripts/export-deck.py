#!/usr/bin/env python
"""
export-deck.py — eksport decku CloudBots do PPTX, PDF i PNG.

Robi to samo co przyciski „PPTX" i „PDF" w decku, tylko lepiej: renderuje
Chromium bez sieci (poza fontami), w pełnej rozdzielczości i bez ograniczeń
html2canvas — gradienty stożkowe, clip-path i maski wychodzą wiernie.

    python export-deck.py                          # PPTX + PDF, ×2
    python export-deck.py --file oferta.html
    python export-deck.py --only pptx              # pptx | pdf | png
    python export-deck.py --scale 3 --name Oferta-Kingfisher
    python export-deck.py --slides 1 4 7           # tylko wskazane slajdy

Slajdy zapisują się jako obrazy — wiernie, ale nieedytowalnie. Gdy klient
potrzebuje PowerPointa z tekstem do poprawek, odbuduj slajdy skilem `pptx`.

Wymaga:
    pip install playwright python-pptx pillow
    playwright install chromium
(python-pptx tylko dla PPTX, pillow tylko dla PDF.)
"""
import argparse
import glob
import http.server
import os
import socketserver
import sys
import threading

from playwright.sync_api import sync_playwright

# konsola Windows domyślnie cp1252 — polskie znaki i strzałki w logu ją wywracają
try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass

PORT = 8813
ROOT = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(ROOT, "_export")

# 13,333 × 7,5 cala to 16:9 w PowerPoincie. Slajd wchodzi na całe płótno,
# bez marginesu — deck jest zaprojektowany jako pełny kadr, nie jako treść w ramce.
PPT_W, PPT_H = 13.333, 7.5

# Chowa sterowanie (topbar, menu, strzałki, pasek postępu) tą samą klasą,
# której używa wbudowany eksport — dzięki temu zrzut wygląda jak slajd,
# a nie jak zrzut ekranu z przeglądarki.
PREPARE = """() => {
  document.body.classList.add('is-clean');
  document.body.classList.remove('drawer-open');
  window.dispatchEvent(new Event('resize'));   /* wymuś przeliczenie autoskalowania */
}"""


def serve(directory):
    handler = lambda *a, **kw: http.server.SimpleHTTPRequestHandler(*a, directory=directory, **kw)
    socketserver.TCPServer.allow_reuse_address = True
    httpd = socketserver.TCPServer(("127.0.0.1", PORT), handler)
    threading.Thread(target=httpd.serve_forever, daemon=True).start()
    return httpd


def pick_file(explicit):
    if explicit:
        return explicit
    found = sorted(glob.glob("*deck*.html")) or sorted(glob.glob("*.html"))
    if not found:
        sys.exit("Nie znalazłem pliku HTML w " + os.getcwd())
    return found[0]


def shoot(target, out_dir, width, height, scale, wanted):
    """Zwraca listę ścieżek PNG w kolejności slajdów."""
    os.makedirs(out_dir, exist_ok=True)
    httpd = serve(os.getcwd())
    paths = []
    try:
        with sync_playwright() as p:
            browser = p.chromium.launch()
            page = browser.new_page(
                viewport={"width": width, "height": height},
                device_scale_factor=scale,
            )
            page.goto(f"http://127.0.0.1:{PORT}/{target}")
            page.wait_for_load_state("networkidle")
            page.evaluate("() => document.fonts.ready")
            page.evaluate(PREPARE)
            page.wait_for_timeout(1200)  # fonty + dopasowanie treści do slajdu

            count = page.locator(".slide").count()
            if not count:
                sys.exit("W pliku nie ma elementów .slide — to nie jest deck")
            nums = wanted or range(1, count + 1)
            print(f"slajdów: {count}   zrzut: {width}×{height} ×{scale}")

            for n in nums:
                if not 1 <= n <= count:
                    print(f"  pomijam {n} — poza zakresem")
                    continue
                page.evaluate(
                    "i => document.querySelectorAll('.slide')[i]"
                    ".scrollIntoView({behavior:'instant',block:'start'})",
                    n - 1,
                )
                page.wait_for_timeout(600)   # animacja wejścia bloków slajdu
                path = os.path.join(out_dir, f"slide-{n:02d}.png")
                page.locator(".slide").nth(n - 1).screenshot(path=path)
                print("  →", path)
                paths.append(path)

            browser.close()
    finally:
        httpd.shutdown()
    return paths


def to_pptx(paths, target_file):
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        sys.exit("PPTX wymaga python-pptx:  pip install python-pptx")
    prs = Presentation()
    prs.slide_width = Inches(PPT_W)
    prs.slide_height = Inches(PPT_H)
    blank = prs.slide_layouts[6]          # układ bez placeholderów
    for png in paths:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(png, 0, 0, width=Inches(PPT_W), height=Inches(PPT_H))
    prs.save(target_file)
    print(f"  → {target_file}  ({len(paths)} slajdów)")


def to_pdf(paths, target_file):
    try:
        from PIL import Image
    except ImportError:
        sys.exit("PDF wymaga Pillow:  pip install pillow")
    pages = [Image.open(p).convert("RGB") for p in paths]
    pages[0].save(target_file, save_all=True, append_images=pages[1:], resolution=150.0)
    print(f"  → {target_file}  ({len(pages)} stron)")


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--file", default=None, help="plik HTML decku")
    ap.add_argument("--out", default=OUT, help="katalog wyjściowy")
    ap.add_argument("--name", default="CloudBots-prezentacja", help="nazwa plików wynikowych")
    ap.add_argument("--scale", type=int, default=2, help="mnożnik rozdzielczości (domyślnie 2)")
    ap.add_argument("--width", type=int, default=1920)
    ap.add_argument("--height", type=int, default=1080)
    ap.add_argument("--slides", nargs="*", type=int, default=None, help="numery slajdów od 1")
    ap.add_argument("--only", choices=["pptx", "pdf", "png"], default=None,
                    help="tylko jeden format (domyślnie PPTX i PDF)")
    args = ap.parse_args()

    target = pick_file(args.file)
    print(f"deck: {target}")

    paths = shoot(target, args.out, args.width, args.height, args.scale, args.slides)
    if not paths:
        sys.exit("Nic nie wyrenderowano")

    if args.only in (None, "pptx"):
        to_pptx(paths, os.path.join(args.out, args.name + ".pptx"))
    if args.only in (None, "pdf"):
        to_pdf(paths, os.path.join(args.out, args.name + ".pdf"))
    if args.only == "png":
        print("  PNG zostały w " + args.out)


if __name__ == "__main__":
    main()
